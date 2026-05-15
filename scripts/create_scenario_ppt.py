#!/usr/bin/env python3
"""
生成客户调研破冰 PPT - 具体场景 before/after 对比
"""

from pptx import Presentation
from pptx.util import Pt
import sys

def create_scenario_ppt(customer_name, output_path):
    """创建场景化破冰 PPT - before/after 对比"""
    
    prs = Presentation()
    
    # 定义具体业务场景（美中宜和 - 妇儿医院场景）
    slides_content = [
        {
            "title": f"{customer_name}\n业务场景洞察",
            "subtitle": "从数据孤岛到智能决策的真实改变",
            "type": "title"
        },
        {
            "title": "场景 1：院长经营看板",
            "before": [
                "每月 5 号等财务报表，各院区 Excel 汇总要 3 天",
                "产科/妇科/儿科收入占比靠手工计算",
                "发现某院区成本异常时，已经是下月中旬",
                "董事会问'哪个科室增长最快'，当场答不上来"
            ],
            "after": [
                "每天早上 8 点自动推送昨日经营日报",
                "实时查看各院区、各科室收入/成本/利润",
                "成本异常自动预警，当天发现当天处理",
                "随时调取任意维度数据，决策有依据"
            ],
            "value": "决策时效：3 天 → 实时；数据准确性：90% → 100%"
        },
        {
            "title": "场景 2：产科护士长排班",
            "before": [
                "每周花 4 小时排班，要平衡 20+ 护士的班次",
                "临时请假找不到人，只能自己顶班",
                "忙闲不均：周一产房忙死，周三护士闲置",
                "护士抱怨'总是我上夜班'，满意度下降"
            ],
            "after": [
                "系统根据历史分娩量预测每日工作量",
                "自动排班 + 一键调整，30 分钟完成",
                "实时显示各时段人力缺口，提前调配",
                "夜班自动轮转，公平透明，护士满意度提升"
            ],
            "value": "排班时间：4 小时 → 30 分钟；护士满意度：65% → 85%"
        },
        {
            "title": "场景 3：患者随访管理",
            "before": [
                "产后 42 天复查靠护士打电话，100 个患者打半天",
                "30% 患者失访，错过复查和二次营销机会",
                "患者问'我上次检查指标怎么样'，要翻病历",
                "投诉处理靠 Excel 记录，无法分析趋势"
            ],
            "after": [
                "系统自动推送复查提醒（微信/短信）",
                "随访率提升到 85%+，二次转化率提升 20%",
                "患者画像完整显示：历史检查、消费记录、偏好",
                "投诉自动分类统计，发现服务薄弱环节"
            ],
            "value": "随访率：70% → 85%；二次转化：10% → 20%"
        },
        {
            "title": "场景 4：库存与耗材管理",
            "before": [
                "药房/耗材库每周盘点，发现缺货已经影响手术",
                "高值耗材（如吻合器）过期报废，一年损失 10 万+",
                "采购凭经验，旺季不够用、淡季堆仓库",
                "财务对账时发现'账实不符'，查不清楚"
            ],
            "after": [
                "库存低于安全线自动预警，提前采购",
                "近效期耗材提前 3 个月预警，优先使用",
                "根据历史消耗 + 预约手术量智能预测采购",
                "出入库自动同步财务，账实实时一致"
            ],
            "value": "缺货率：15% → 2%；耗材报废：10 万/年 → 1 万/年"
        },
        {
            "title": "下一步验证计划",
            "points": [
                "第 1 周：选 1-2 个痛点场景深度调研（如院长看板 + 排班）",
                "第 2-3 周：用真实数据搭建 POC 环境",
                "第 4 周：业务部门试用，量化价值（时间节省/错误减少）",
                "第 5 周：基于 POC 结果决策是否规模化推广"
            ]
        }
    ]
    
    # 创建幻灯片
    for i, slide_data in enumerate(slides_content):
        if slide_data.get("type") == "title":
            # 封面页
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
        elif "before" in slide_data:
            # Before/After 对比页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            tf = content.text_frame
            tf.clear()
            
            # Before 部分
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "❌ 现在（Before）"
            run.font.bold = True
            run.font.size = Pt(20)
            
            for point in slide_data["before"]:
                p = tf.add_paragraph()
                p.text = f"  • {point}"
                p.font.size = Pt(16)
                p.level = 0
                p.space_after = Pt(8)
            
            # 空一行
            p = tf.add_paragraph()
            p.text = ""
            
            # After 部分
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "✅ 未来（After）"
            run.font.bold = True
            run.font.size = Pt(20)
            
            for point in slide_data["after"]:
                p = tf.add_paragraph()
                p.text = f"  • {point}"
                p.font.size = Pt(16)
                p.level = 0
                p.space_after = Pt(8)
            
            # 价值总结
            p = tf.add_paragraph()
            p.text = ""
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"📊 价值：{slide_data.get('value', '')}"
            run.font.bold = True
            run.font.size = Pt(18)
        else:
            # 普通内容页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            tf = content.text_frame
            tf.clear()
            
            for j, point in enumerate(slide_data["points"], 1):
                if j == 1:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = point
                p.font.size = Pt(18)
                p.level = 0
                p.space_after = Pt(12)
    
    # 保存 PPT
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法：python create_scenario_ppt.py <客户名称> <输出路径>")
        sys.exit(1)
    
    customer_name = sys.argv[1]
    output_path = sys.argv[2]
    
    create_scenario_ppt(customer_name, output_path)
