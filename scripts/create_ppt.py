#!/usr/bin/env python3
"""
生成客户调研破冰 PPT - 3-5 页客户视角场景价值
"""

from pptx import Presentation
from pptx.util import Pt
import sys

def create_icebreaker_ppt(customer_name, output_path):
    """创建破冰 PPT"""
    
    prs = Presentation()
    
    # 定义客户视角的场景价值内容（不谈技术）
    slides_content = [
        {
            "title": f"{customer_name}\n数字化场景洞察",
            "subtitle": "业务价值驱动的增长路径",
            "type": "title"
        },
        {
            "title": "行业趋势与标杆实践",
            "points": [
                "医疗行业数字化转型进入深水区",
                "AI+ 数据驱动成为头部医院标配",
                "政策持续鼓励智慧医院建设",
                "患者体验与运营效率双提升成核心竞争力"
            ]
        },
        {
            "title": f"{customer_name} 业务场景分析",
            "points": [
                "多院区协同：数据实时互通，管理决策统一",
                "临床支持：诊疗数据沉淀，辅助医学决策",
                "运营优化：资源动态调度，成本精细管控",
                "患者服务：全流程体验优化，满意度持续提升"
            ]
        },
        {
            "title": "可量化的业务价值",
            "points": [
                "管理决策效率提升 50%+",
                "数据报表时间从小时级降至分钟级",
                "跨院区数据同步延迟<1 分钟",
                "运营管理成本降低 30%+",
                "患者等待时间缩短 40%+"
            ]
        },
        {
            "title": "建议验证路径",
            "points": [
                "第 1 周：业务场景深度调研",
                "第 2-3 周：POC 环境搭建与场景配置",
                "第 4 周：业务价值验证与评估",
                "第 5 周：规模化推广规划"
            ]
        }
    ]
    
    # 创建幻灯片
    for i, slide_data in enumerate(slides_content):
        if slide_data.get("type") == "title":
            # 封面页
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
        else:
            # 内容页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            # 添加要点
            tf = content.text_frame
            tf.clear()
            
            for j, point in enumerate(slide_data["points"]):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = point
                p.font.size = Pt(18)
                p.level = 0
                p.space_after = Pt(12)
    
    # 保存 PPT
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法：python create_ppt.py <客户名称> <输出路径>")
        sys.exit(1)
    
    customer_name = sys.argv[1]
    output_path = sys.argv[2]
    
    create_icebreaker_ppt(customer_name, output_path)
